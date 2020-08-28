from pptx import Presentation
from pptx.util import Inches
import io
from os.path import join as opj
import os

__all__ = ['BasicPresentation']

""" (text, fig, table) """
layouts = {(True, True, True):0,
           (True, True, False):1,
           (True, False, True):2,
           (True, False, False):3,
           (False, True, False):4,
           (False, False, True):5}

"""TODO:
    - Add title slide with date, name and title"""

class BasicPresentation():
    def __init__(self, filename):
        """Create a basic slide deck.
        Note: Figures should be
              
              6.5" x 6" (W x H) w/o a table
              6.5" x 3.5" w/ table
              12.5" x 6" alone

              or will need resizing"""
        self.filename = filename
        try:
            self.prs = Presentation(opj(os.path.dirname(__file__), 'basic_pptx.pptx'))
        except NameError:
            self.prs = Presentation()

    def add_slide(self, title=None, text=None, figure=None, table=None):
        lo = (not text is None, not figure is None, not table is None)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layouts[lo]])
        
        if not title is None:
            title_shape = slide.shapes.title.text = title
        
        if not text is None:
            body_shape = slide.shapes.placeholders[1].text_frame.text = text[0]
            for t in text[1:]:
                p = slide.shapes.placeholders[1].text_frame.add_paragraph()
                p.text = t
                p.level = 0

        if not figure is None:
            streamio = io.BytesIO()
            figure.savefig(streamio, format='png', dpi=200)
            streamio.seek(0)
            # pic = slide.shapes.add_picture(streamio, Inches(5), Inches(0.5))
            pic = slide.shapes.placeholders[10].insert_picture(streamio)
        if not table is None:
            tab = slide.shapes.placeholders[11].insert_table(table.shape[0] + 1, table.shape[1]).table
            """tab = slide.shapes.add_table(table.shape[0] + 1, table.shape[1],
                                         left=Inches(1),
                                         top=Inches(4),
                                         width=Inches(6),
                                         height=Inches(4)).table"""
            for j in range(table.shape[1]):
                if type(table.columns[j]) is tuple:
                    col = '|'.join([str(c) for c in table.columns[j]])
                else:
                    col = str(table.columns[j])
                tab.cell(0, j).text = col
            
            for i in range(table.shape[0]):
                for j in range(table.shape[1]):
                    tab.cell(i + 1, j).text = str(table.values[i, j])
    def save(self):
        self.prs.save(self.filename)



def test_basic():
    import pandas as pd
    import matplotlib.pyplot as plt
    import numpy as np
    from fg_shared import _git

    df = pd.DataFrame({'dude':np.arange(5), 'yo':np.linspace(4, 5, 5)})

    bp = BasicPresentation(opj(_git, 'test.pptx'))

    figh = plt.figure(figsize=(13, 7))
    plt.plot(np.random.randn(50), np.random.rand(50), 'o', alpha=0.5)
    bp.add_slide(title='Dude, title!',
                 text=['Bullet 1', 'Bullet 2', 'Bullet 3'],
                 table=df,
                 figure=figh)

    bp.add_slide(title='Dude, title!',
                 text=['Bullet 1', 'Bullet 2', 'Bullet 3'])
    bp.add_slide(title='Dude, title!',
                 text=['Bullet 1', 'Bullet 2', 'Bullet 3'],
                 table=df)
    
    figh = plt.figure(figsize=(6.5, 6))
    plt.plot(np.random.randn(50), np.random.rand(50), 'o', alpha=0.5)
    bp.add_slide(title='Dude, title!',
                 text=['Bullet 1', 'Bullet 2', 'Bullet 3'],
                 figure=figh)
    
    figh = plt.figure(figsize=(12.5, 6))
    plt.plot(np.random.randn(50), np.random.rand(50), 'o', alpha=0.5)
    bp.add_slide(title='Dude, title!',
                 figure=figh)
    bp.add_slide(title='Dude, title!',
                 table=df)
    bp.save()



